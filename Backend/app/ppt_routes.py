from flask import Blueprint, request, send_file
from PyPDF2 import PdfMerger
import io
from pdf2docx import Converter
import os
import tempfile
import subprocess
from docx import Document
from PIL import Image
from fpdf import FPDF
from pdf2image import convert_from_path
import zipfile
from pptx import Presentation
from pptx.util import Inches


ppt_bp = Blueprint('ppt', __name__)


@ppt_bp.route('/convert_ppt_to_pdf', methods=['POST'])
def convert_ppt_to_pdf():
    if 'ppt' not in request.files:
        return "No PowerPoint file part in the request", 400

    ppt_file = request.files['ppt']
    
    if not ppt_file.filename.endswith(('.ppt', '.pptx')):
        return "Please upload a valid PowerPoint file", 400

    # Create an in-memory file to save the uploaded PowerPoint file
    ppt_file_bytes = io.BytesIO(ppt_file.read())

    # Use a temporary directory for saving converted PDF, or an in-memory solution
    temp_pdf_path = io.BytesIO()

    try:
        # Save ppt file to a temporary path in memory for conversion
        ppt_file_path = 'temp_ppt.pptx'
        with open(ppt_file_path, 'wb') as temp_file:
            temp_file.write(ppt_file_bytes.getbuffer())
        
        # Convert the PowerPoint file to PDF using LibreOffice
        subprocess.run([r'C:\Program Files\LibreOffice\program\soffice.exe',
                        '--headless', '--convert-to', 'pdf', ppt_file_path,
                        '--outdir', os.getcwd()], check=True)

        # Open the generated PDF file and store it in the in-memory buffer
        with open(ppt_file_path.rsplit('.', 1)[0] + '.pdf', 'rb') as pdf_file:
            temp_pdf_path.write(pdf_file.read())
        
        # Close the in-memory buffer and the temp ppt file
        temp_pdf_path.seek(0)

        # Send the PDF file back to the client
        return send_file(temp_pdf_path, as_attachment=True,
                         download_name="converted.pdf",
                         mimetype="application/pdf")

    except subprocess.CalledProcessError as e:
        return f"Error converting PowerPoint to PDF: {e}", 500

    finally:
        if os.path.exists(ppt_file_path):
            os.remove(ppt_file_path) 

@ppt_bp.route('/convert_pdf_to_ppt', methods=['POST'])
def convert_pdf_to_ppt():
    if 'pdf' not in request.files:
        return "No PDF file part in the request", 400

    pdf_file = request.files['pdf']
    
    if not pdf_file.filename.endswith('.pdf'):
        return "Please upload a valid PDF file", 400

    # Convert PDF to images (one image per page)
    pdf_file_bytes = io.BytesIO(pdf_file.read())
    
    # Save the PDF byte stream to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf_file:
        temp_pdf_file.write(pdf_file_bytes.read())
        temp_pdf_path = temp_pdf_file.name

    # Convert the saved PDF file to images
    images = convert_from_path(temp_pdf_path, 300)  # 300 DPI is high-quality for presentation

    # Create a PowerPoint presentation
    ppt = Presentation()

    # Add images as slides
    for img in images:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Using a blank slide layout
        img_stream = io.BytesIO()
        img.save(img_stream, format='PNG')  # Save image to the in-memory stream
        img_stream.seek(0)

        # Add the image to the slide
        slide.shapes.add_picture(img_stream, 
                                 Inches(0), Inches(0), 
                                 width=Inches(10), 
                                 height=Inches(7.5))  # Adjust size as needed

    # Save the PowerPoint file to an in-memory buffer
    ppt_output = io.BytesIO()
    ppt.save(ppt_output)
    ppt_output.seek(0)

    # Clean up the temporary PDF file
    os.remove(temp_pdf_path)

    # Send the PowerPoint file as a download
    return send_file(ppt_output, as_attachment=True,
                     download_name="converted.pptx",
                     mimetype=
                     "application/vnd.openxmlformats-officedocument.presentationml.presentation")



