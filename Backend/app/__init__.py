from flask import Flask, request, send_file
from PyPDF2 import PdfMerger
import io
from pdf2docx import Converter
import os
import tempfile
import subprocess
from docx import Document
from PIL import Image
from fpdf import FPDF
from pdf2image import convert_from_path
import zipfile
from pptx import Presentation
from pptx.util import Inches

app = Flask(__name__)

@app.route('/merge_pdfs', methods=['POST'])
def merge_pdfs():
    if 'pdfs' not in request.files:
        return "No files part in the request", 400

    files = request.files.getlist('pdfs')
    
    if len(files) < 2:
        return "Please upload at least two PDF files to merge", 400

    merger = PdfMerger()

    for file in files:
        if file.filename.endswith('.pdf'):
            merger.append(file)
        else:
            return f"File {file.filename} is not a PDF", 400

    merged_pdf = io.BytesIO()
    merger.write(merged_pdf)
    merger.close()
    
    merged_pdf.seek(0)

    return send_file(merged_pdf, as_attachment=True, download_name="merged.pdf", mimetype="application/pdf")

@app.route('/convert_pdf_to_doc', methods=['POST'])
def convert_pdf_to_doc():
    if 'pdf' not in request.files:
        return "No file part", 400

    file = request.files['pdf']

    if file.filename == '':
        return "No selected file", 400

    if file and file.filename.endswith('.pdf'):
        pdf_bytes = file.read()
        
        pdf_io = io.BytesIO(pdf_bytes)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf:
            temp_pdf.write(pdf_bytes)
            temp_pdf_path = temp_pdf.name

        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_docx:
            docx_path = temp_docx.name

        converter = Converter(temp_pdf_path)
        converter.convert(docx_path, start=0, end=None)
# MIME Type: It tells the browser what kind of file is being sent so it knows how to handle it.
# For DOCX files, the MIME type
# is application/vnd.openxmlformats-officedocument.wordprocessingml.document
        return send_file(docx_path, as_attachment=True, download_name="converted.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    return "Invalid file format", 400


@app.route('/convert_doc_to_pdf', methods=['POST'])
def convert_doc_to_pdf():
    if 'doc' not in request.files:
        return "No file part", 400

    file = request.files['doc']

    if file.filename == '':
        return "No selected file", 400

    if file and file.filename.endswith('.docx'):
        doc = file.read()
        docx_data = io.BytesIO(doc)

        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_docx:
            temp_docx.write(doc)
            temp_docx_path = temp_docx.name

        temp_pdf_path = temp_docx_path.replace('.docx', '.pdf')

        try:
            subprocess.run([r'C:\Program Files\LibreOffice\program\soffice.exe', 
                            '--headless', '--convert-to', 'pdf', temp_docx_path, '--outdir',
                            os.path.dirname(temp_pdf_path)], check=True)

            with open(temp_pdf_path, 'rb') as pdf_file:
                pdf_data = io.BytesIO(pdf_file.read())

            pdf_data.seek(0)
            return send_file(pdf_data, as_attachment=True, download_name="converted.pdf", mimetype="application/pdf")

        finally:
            os.remove(temp_docx_path)
            if os.path.exists(temp_pdf_path):
                os.remove(temp_pdf_path)

    return "Invalid file format", 400

@app.route('/merge_docs', methods=['POST'])
def merge_docs():
    if 'docs' not in request.files:
        return "No files part in the request", 400

    files = request.files.getlist('docs')

    if len(files) < 2:
        return "Please upload at least two DOCX files to merge", 400

    merged_doc = Document()

    for file in files:
        if file.filename.endswith('.docx'):
            doc = Document(file)
            
            for element in doc.element.body:
                merged_doc.element.body.append(element)
        else:
            return f"File {file.filename} is not a DOCX", 400

    merged_doc_io = io.BytesIO()
    merged_doc.save(merged_doc_io)
    merged_doc_io.seek(0)

    return send_file(merged_doc_io, as_attachment=True, download_name="merged_document.docx", mimetype="application/vnd.openxmlformats-officedocument.wordprocessingml.document")

@app.route('/convert_images_to_pdf', methods=['POST'])
def convert_images_to_pdf():
    print("Request files:", request.files)  # Log all incoming files
    if 'images' not in request.files:
        return "No 'images' field in the request", 400

    files = request.files.getlist('images')
    
    if not files:
        return "Please upload at least one image", 400

    image_list = []
    for file in files:
        image = Image.open(file)
        image_list.append(image)

    pdf_output = io.BytesIO()
    image_list[0].save(pdf_output, "PDF", save_all=True, append_images=image_list[1:])
    
    pdf_output.seek(0)
    return send_file(pdf_output, as_attachment=True, download_name="converted_images.pdf", mimetype="application/pdf")

@app.route('/convert_pdf_to_images', methods=['POST'])
def convert_pdf_to_images():
    if 'pdf' not in request.files:
        return "No PDF file part", 400

    file = request.files['pdf']
    
    if file.filename == '':
        return "No selected file", 400

    if file and file.filename.endswith('.pdf'):
        # Read the uploaded PDF file into a BytesIO object
        pdf_data = file.read()
        
        # Create a temporary file to save the uploaded PDF
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_pdf:
            tmp_pdf.write(pdf_data)
            tmp_pdf_path = tmp_pdf.name  # Store the path of the temporary file

        # Convert PDF to images (one image per page)
        try:
            images = convert_from_path(tmp_pdf_path)
            
            # Prepare a zip file to store the images
            zip_filename = "pdf_images.zip"
            zip_buffer = io.BytesIO()

            with zipfile.ZipFile(zip_buffer, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for i, image in enumerate(images):
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format='PNG')
                    img_byte_arr.seek(0)
                    zipf.writestr(f"page_{i + 1}.png", img_byte_arr.read())

            zip_buffer.seek(0)
            return send_file(zip_buffer, as_attachment=True, download_name=zip_filename, mimetype="application/zip")
        
        except Exception as e:
            return f"Error converting PDF: {str(e)}", 500
        
        finally:
            # Clean up the temporary PDF file after conversion
            if os.path.exists(tmp_pdf_path):
                os.remove(tmp_pdf_path)

    return "Invalid file format", 400

@app.route('/convert_ppt_to_pdf', methods=['POST'])
def convert_ppt_to_pdf():
    if 'ppt' not in request.files:
        return "No PowerPoint file part in the request", 400

    ppt_file = request.files['ppt']
    
    if not ppt_file.filename.endswith(('.ppt', '.pptx')):
        return "Please upload a valid PowerPoint file", 400

    # Create an in-memory file to save the uploaded PowerPoint file
    ppt_file_bytes = io.BytesIO(ppt_file.read())

    # Use a temporary directory for saving converted PDF, or an in-memory solution
    temp_pdf_path = io.BytesIO()

    try:
        # Save ppt file to a temporary path in memory for conversion
        ppt_file_path = 'temp_ppt.pptx'
        with open(ppt_file_path, 'wb') as temp_file:
            temp_file.write(ppt_file_bytes.getbuffer())
        
        # Convert the PowerPoint file to PDF using LibreOffice
        subprocess.run([r'C:\Program Files\LibreOffice\program\soffice.exe',
                        '--headless', '--convert-to', 'pdf', ppt_file_path,
                        '--outdir', os.getcwd()], check=True)

        # Open the generated PDF file and store it in the in-memory buffer
        with open(ppt_file_path.rsplit('.', 1)[0] + '.pdf', 'rb') as pdf_file:
            temp_pdf_path.write(pdf_file.read())
        
        # Close the in-memory buffer and the temp ppt file
        temp_pdf_path.seek(0)

        # Send the PDF file back to the client
        return send_file(temp_pdf_path, as_attachment=True, download_name="converted.pdf", mimetype="application/pdf")

    except subprocess.CalledProcessError as e:
        return f"Error converting PowerPoint to PDF: {e}", 500

    finally:
        if os.path.exists(ppt_file_path):
            os.remove(ppt_file_path) 

@app.route('/convert_pdf_to_ppt', methods=['POST'])
def convert_pdf_to_ppt():
    if 'pdf' not in request.files:
        return "No PDF file part in the request", 400

    pdf_file = request.files['pdf']
    
    if not pdf_file.filename.endswith('.pdf'):
        return "Please upload a valid PDF file", 400

    # Convert PDF to images (one image per page)
    pdf_file_bytes = io.BytesIO(pdf_file.read())
    
    # Save the PDF byte stream to a temporary file
    with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_pdf_file:
        temp_pdf_file.write(pdf_file_bytes.read())
        temp_pdf_path = temp_pdf_file.name

    # Convert the saved PDF file to images
    images = convert_from_path(temp_pdf_path, 300)  # 300 DPI is high-quality for presentation

    # Create a PowerPoint presentation
    ppt = Presentation()

    # Add images as slides
    for img in images:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])  # Using a blank slide layout
        img_stream = io.BytesIO()
        img.save(img_stream, format='PNG')  # Save image to the in-memory stream
        img_stream.seek(0)

        # Add the image to the slide
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))  # Adjust size as needed

    # Save the PowerPoint file to an in-memory buffer
    ppt_output = io.BytesIO()
    ppt.save(ppt_output)
    ppt_output.seek(0)

    # Clean up the temporary PDF file
    os.remove(temp_pdf_path)

    # Send the PowerPoint file as a download
    return send_file(ppt_output, as_attachment=True, download_name="converted.pptx", mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation")






if __name__ == '__main__':
    app.run(debug=True)

